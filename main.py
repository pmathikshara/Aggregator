from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from docx import Document

import os
from pathlib import Path
import glob2
import re

for eachfile in glob2.glob("*.pptx"):
    prs = Presentation(eachfile)
    print(eachfile)
    print("----------------------")
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                # print(shape.text)
                print(re.search("[\$\£\€](\d+(?:\.\d{1,2})?)", shape.text))

# for x in range(1,16,1):
#     num = "0" + str(x)
#     num_pad = num[-2:]
#     filename = "{}_numbered_items.pptx".format(num_pad, num_pad)
#     print(filename)

#     check_file = Path(filename)
#     if check_file.is_file():
#         prs = Presentation(filename)
#         slides = prs.slides
#         image_counter = 0
#         run_text = []
#         for slide in prs.slides:
#             for shape in slide.shapes:
#                 if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
#                     image = shape.image
#                     # ---get image "file" contents---
#                     image_bytes = image.blob
#                     # ---make up a name for the file, e.g. 'image.jpg'---
#                     image_filename = "{}_lesson_image_{}.{}".format(num_pad, image_counter, image.ext)
#                     with open(image_filename, 'wb') as image_f:
#                         image_f.write(image_bytes)
#                     markdown_image = "![{}](images/{})".format(image_filename, image_filename)
#                     run_text.append(markdown_image)
#                     image_counter = image_counter + 1
#                 if not shape.has_text_frame:
#                     continue
#                 for paragraph in shape.text_frame.paragraphs:
#                     run_text.append(paragraph.text)
#                     # for run in paragraph.runs:
#                     #     print(run.text)

#             if slide.has_notes_slide:
#                 notes_slide = slide.notes_slide
#                 for paragraph in notes_slide.notes_text_frame.paragraphs:
#                     run_text.append(paragraph.text)
#                 # notes_slide.notes_text_frame

#         lesson_output_filename = "{}_lesson_out.txt".format(num_pad)
#         with open(lesson_output_filename, 'w') as output_f:
#             for item in run_text:
#                 try:
#                     output_f.write("{}\n".format(item))
#                 except UnicodeEncodeError as ex:
#                     print("slide file {}".format(num_pad))
#                     print(item)

# ## and now DOCX
# import docx2txt

# docx_files_by_hand = [
#     "manually.docx",
# ]

# docx_files_by_docx2txt = [
#     "automated.docx",
# ]

# for filename in docx_files_by_hand:
#     print(filename)
#     check_file = Path(filename)
#     if check_file.is_file():
#         docu = Document(filename)
#         basename = filename.split(".")
#         run_text = []
#         for paragraph in docu.paragraphs:
#             run_text.append(paragraph.text)
#         lesson_output_filename = "{}_trans.txt".format(basename[0])
#         with open(lesson_output_filename, 'w') as output_f:
#             for item in run_text:
#                 try:
#                     output_f.write("{}\n".format(item))
#                 except UnicodeEncodeError as ex:
#                     print(ex)

# for filename in docx_files_by_docx2txt:
#     print(filename)

#     check_file = Path(filename)
#     if check_file.is_file():
#         basename = filename.split(".")

#         # Create target Directory if don't exist
#         if not os.path.exists(basename[0]):
#             os.mkdir(basename[0])
#             print("Directory ", basename[0], " Created ")
#         else:
#             print("Directory ", basename[0], " already exists")

#         text = docx2txt.process(filename, basename[0])
#         lesson_output_filename = "{}_trans.txt".format(basename[0])
#         with open(lesson_output_filename, 'w') as output_f:
#             try:
#                 output_f.write("{}".format(text))
#             except UnicodeEncodeError as ex:
#                 print("doc file {}".format(basename[0]))
